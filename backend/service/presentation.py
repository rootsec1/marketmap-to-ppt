from pptx import Presentation
from pptx.util import Inches
import requests
from io import BytesIO
import os


def create_ppt_with_logos(logo_sources, is_local=True, slide_width=10, slide_height=7.5, margin=0.5):
    """
    Create a PowerPoint presentation with logos evenly arranged in a grid using the entire slide space.

    :param logo_sources: List of local file paths or URLs.
    :param is_local: Boolean indicating whether logos are from local files or remote URLs.
    :param slide_width: Width of the slide in inches (default 10).
    :param slide_height: Height of the slide in inches (default 7.5).
    :param margin: Margin around the slide in inches (default 0.5).
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Calculate available space
    available_width = slide_width - 2 * margin
    available_height = slide_height - 2 * margin

    # Determine grid size
    num_images = len(logo_sources)
    grid_columns = int(num_images ** 0.5)  # Approximate square grid
    if grid_columns == 0:
        grid_columns = 1
    grid_rows = (num_images + grid_columns - 1) // grid_columns  # Ceiling division

    # Adjust logo size based on grid
    logo_width = available_width / grid_columns
    logo_height = available_height / grid_rows

    # Starting position
    x_start = Inches(margin)
    y_start = Inches(margin)

    # Loop to place each logo
    x, y = x_start, y_start
    current_column = 0

    for source in logo_sources:
        try:
            # Add the logo to the slide
            if is_local:
                if os.path.exists(source):
                    slide.shapes.add_picture(source, x, y, width=Inches(logo_width), height=Inches(logo_height))
                else:
                    print(f"⚠️ File not found: {source}, skipping...")
            else:
                response = requests.get(source)
                if response.status_code == 200:
                    image_stream = BytesIO(response.content)
                    slide.shapes.add_picture(image_stream, x, y, width=Inches(logo_width), height=Inches(logo_height))
                else:
                    print(f"⚠️ Failed to fetch image from: {source}, skipping...")

            # Update position for next logo
            current_column += 1
            if current_column >= grid_columns:  # Move to the next row
                current_column = 0
                x = x_start
                y += Inches(logo_height)
            else:
                x += Inches(logo_width)

        except Exception as e:
            print(f"❌ Error processing {source}: {str(e)}")

    # Save the PowerPoint file
    ppt_path = "tmp/presentation.pptx"
    prs.save(ppt_path)
    print(f"✅ PowerPoint file saved as {ppt_path}")

    return ppt_path
